#!/usr/bin/env python3
"""
=================================================================
  PII DATA DISCOVERY TOOL v3.0 — India Regulatory Focus
  Fully Offline | ML-Enhanced | Industry-Standard Reporting
=================================================================

Detection Pipeline (6 Layers):
  1. Regex Pattern Matching  — Structured Indian PII formats
  2. spaCy NER (en_core_web_lg) — Person names, organisations, locations
  3. Column Header Analysis   — Excel/CSV schema PII detection
  4. Checksum Validation      — Luhn (cards), Verhoeff (Aadhaar)
  5. Context Scoring          — Confidence scoring per finding
  6. Image/OCR Detection      — Embedded images in DOCX/PDF

Applicable Laws & Regulators:
  - Digital Personal Data Protection Act, 2023 (DPDP Act)
  - DPDP Rules, 2025 (notified Nov 2025)
  - Information Technology Act, 2000 (IT Act)
  - IT (Reasonable Security Practices) Rules, 2011
  - RBI Master Directions (Data Localisation, Card-on-File, IT Governance)
  - SEBI Cybersecurity & Cyber Resilience Framework (CSCRF) 2023
  - IRDAI Information & Cybersecurity Guidelines 2023
  - PFRDA Subscriber Data Protection Guidelines
  - TRAI Telecom Subscriber Privacy Regulations
  - PCI-DSS v4.0
  - CERT-In Directions 2022 (6-hour mandatory incident reporting)
  - Aadhaar (Targeted Delivery) Act, 2016

Usage:
  python pii_scanner_india.py
  python pii_scanner_india.py --target "path/to/folder"
  python pii_scanner_india.py --target "path/to/folder" --output "path/to/reports"
"""

from __future__ import annotations

import os
import re
import csv
import bisect
import hashlib
import argparse
import datetime
import pathlib
import zipfile
from collections import defaultdict

# ---------------------------------------------------------------------------
# Third-party imports
# ---------------------------------------------------------------------------
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# spaCy — NER for names, orgs, locations (Layer 2) — Lazy Loader
_NLP = None
_NLP_MODEL_NAME = "none"
_SPACY_INITIALIZED = False
HAS_SPACY = False

def get_nlp():
    global _NLP, _NLP_MODEL_NAME, _SPACY_INITIALIZED, HAS_SPACY
    if _SPACY_INITIALIZED:
        return _NLP, _NLP_MODEL_NAME
    _SPACY_INITIALIZED = True
    disable_components = ["parser", "lemmatizer", "tagger", "attribute_ruler", "morphologizer", "senter"]

    # Method 1: Direct model module load (Primary for PyInstaller standalone binaries & virtual environments)
    import importlib
    for _mod_name in ("en_core_web_lg", "en_core_web_sm"):
        try:
            mod = importlib.import_module(_mod_name)
            _NLP = mod.load(disable=disable_components)
            _NLP_MODEL_NAME = _mod_name
            HAS_SPACY = True
            return _NLP, _NLP_MODEL_NAME
        except Exception:
            continue

    # Method 2: Fallback spacy.load() string lookup
    try:
        import spacy
        for _model_name in ("en_core_web_lg", "en_core_web_sm"):
            try:
                _NLP = spacy.load(_model_name, disable=disable_components)
                _NLP_MODEL_NAME = _model_name
                HAS_SPACY = True
                break
            except OSError:
                continue
    except ImportError:
        pass
    return _NLP, _NLP_MODEL_NAME

# Optional file parsers — degrade gracefully
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# OCR support (optional)
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False


# ===================================================================
#  CONSTANTS
# ===================================================================
BASE_DIR = pathlib.Path(__file__).parent.resolve()
DEFAULT_TARGET = BASE_DIR / "test"
DEFAULT_REPORTS_DIR = BASE_DIR / "reports"
TOOL_VERSION = "3.0.0"
TOOL_NAME = "PII Data Discovery Tool"
REPORT_FILENAME = "PII_Discovery_Report.xlsx"

SCAN_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml",
    ".log", ".sql", ".env", ".ini", ".conf", ".cfg", ".md",
    ".html", ".htm", ".rtf",
    ".pdf",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx",
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif"}


# ===================================================================
#  TEXT SANITISATION
# ===================================================================
_ILLEGAL_XML_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f]')

def _sanitize(text: str | None) -> str:
    if not text:
        return ""
    return _ILLEGAL_XML_RE.sub('', text)


# ===================================================================
#  CHECKSUM VALIDATORS
# ===================================================================

def luhn_check(number_str: str) -> bool:
    digits = [int(d) for d in re.sub(r'\D', '', number_str)]
    if len(digits) < 13 or len(digits) > 19:
        return False
    digits.reverse()
    total = 0
    for i, d in enumerate(digits):
        if i % 2 == 1:
            d *= 2
            if d > 9:
                d -= 9
        total += d
    return total % 10 == 0


_VERHOEFF_D = [
    [0,1,2,3,4,5,6,7,8,9],[1,2,3,4,0,6,7,8,9,5],
    [2,3,4,0,1,7,8,9,5,6],[3,4,0,1,2,8,9,5,6,7],
    [4,0,1,2,3,9,5,6,7,8],[5,9,8,7,6,0,4,3,2,1],
    [6,5,9,8,7,1,0,4,3,2],[7,6,5,9,8,2,1,0,4,3],
    [8,7,6,5,9,3,2,1,0,4],[9,8,7,6,5,4,3,2,1,0]
]
_VERHOEFF_P = [
    [0,1,2,3,4,5,6,7,8,9],[1,5,7,6,2,8,3,0,9,4],
    [5,8,0,3,7,9,6,1,4,2],[8,9,1,6,0,4,3,5,2,7],
    [9,4,5,3,1,2,6,8,7,0],[4,2,8,6,5,7,3,9,0,1],
    [2,7,9,3,8,0,6,4,1,5],[7,0,4,6,9,1,3,2,5,8]
]

def verhoeff_check(number_str: str) -> bool:
    digits = re.sub(r'\D', '', number_str)
    if len(digits) != 12:
        return False
    c = 0
    for i, ch in enumerate(reversed(digits)):
        c = _VERHOEFF_D[c][_VERHOEFF_P[i % 8][int(ch)]]
    return c == 0


# ===================================================================
#  REGEX PATTERN DEFINITIONS (Indian PII)
# ===================================================================

class PIIPattern:
    __slots__ = ('tag', 'regex', 'description', 'sensitivity',
                 'regulation', 'validator', 'min_len')

    def __init__(self, tag, regex, description, sensitivity,
                 regulation, validator=None, min_len=0):
        self.tag = tag
        self.regex = regex
        self.description = description
        self.sensitivity = sensitivity
        self.regulation = regulation
        self.validator = validator
        self.min_len = min_len


def _build_regex_patterns() -> list[PIIPattern]:
    P = PIIPattern
    return [
        P("AADHAAR",
          re.compile(r'\b([2-9]\d{3})\s?(\d{4})\s?(\d{4})\b'),
          "Aadhaar Number (12-digit UIDAI)", "HIGH",
          "DPDP Act 2023 Sec 2(t), Aadhaar Act 2016 Sec 29, DPDP Rules 2025",
          validator=lambda m: verhoeff_check(re.sub(r'\D', '', m.group())),
          min_len=12),

        P("PAN_CARD",
          re.compile(r'\b([A-Z]{5}\d{4}[A-Z])\b'),
          "PAN Card Number", "HIGH",
          "DPDP Act 2023, Income Tax Act Sec 139A, DPDP Rules 2025",
          min_len=10),

        P("PASSPORT_IN",
          re.compile(r'\b([A-PR-WY][1-9]\d{6})\b'),
          "Indian Passport Number", "HIGH",
          "DPDP Act 2023, Passports Act 1967, DPDP Rules 2025",
          min_len=8),

        P("VOTER_ID",
          re.compile(r'\b([A-Z]{3}\d{7})\b'),
          "Voter ID / EPIC Number", "HIGH",
          "DPDP Act 2023, Representation of People Act 1950",
          min_len=10),

        P("DRIVING_LICENSE",
          re.compile(r'\b([A-Z]{2}[-\s]?\d{2}[-\s]?\d{4}[-\s]?\d{7})\b'),
          "Indian Driving License", "HIGH",
          "DPDP Act 2023, Motor Vehicles Act 1988",
          min_len=13),

        P("GSTIN",
          re.compile(r'\b(\d{2}[A-Z]{5}\d{4}[A-Z]\d[Z][A-Z0-9])\b'),
          "GST Identification Number", "MEDIUM",
          "GST Act 2017, DPDP Act 2023",
          min_len=15),

        P("IFSC_CODE",
          re.compile(r'\b([A-Z]{4}0[A-Z0-9]{6})\b'),
          "IFSC Code (Bank Branch)", "MEDIUM",
          "RBI Master Direction on IT Governance 2023, DPDP Act 2023",
          min_len=11),

        P("PHONE_IN",
          re.compile(r'(?<!\d)(\+91[\s\-]?[6-9]\d{9}|0[6-9]\d{9}|[6-9]\d{9})(?!\d)'),
          "Indian Mobile / Phone Number", "HIGH",
          "DPDP Act 2023, TRAI Privacy Regulations, DPDP Rules 2025",
          min_len=10),

        P("EMAIL",
          re.compile(r'\b([A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,})\b'),
          "Email Address", "HIGH",
          "DPDP Act 2023, IT Rules 2011 Sec 3, DPDP Rules 2025",
          min_len=5),

        P("CREDIT_CARD",
          re.compile(r'\b(\d{4}[\s\-]?\d{4}[\s\-]?\d{4}[\s\-]?\d{1,7})\b'),
          "Credit/Debit Card Number", "HIGH",
          "PCI-DSS v4.0, RBI Card-on-File Tokenisation Dir. 2022",
          validator=lambda m: luhn_check(m.group()),
          min_len=13),

        P("UPI_ID",
          re.compile(
              r'\b([A-Za-z0-9._%+\-]+@(?:upi|paytm|ybl|okhdfcbank|oksbi'
              r'|okaxis|okicici|axl|ibl|apl|sbi|icici|hdfc|kotak'
              r'|yesbank|indus|rbl|federal|bob|canara|pnb|boi|cbi'
              r'|union|dbs|idbi|jkb|kvb|cub|scb|axisbank|idfcfirst))\b'),
          "UPI Virtual Payment Address", "MEDIUM",
          "RBI UPI Guidelines, NPCI Circulars, DPDP Act 2023",
          min_len=5),

        P("IP_ADDRESS",
          re.compile(
              r'\b((?:25[0-5]|2[0-4]\d|[01]?\d\d?)\.(?:25[0-5]|2[0-4]\d|[01]?\d\d?)'
              r'\.(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\.(?:25[0-5]|2[0-4]\d|[01]?\d\d?))\b'),
          "IPv4 Address", "LOW",
          "IT Act 2000 Sec 43, CERT-In Directions 2022 (6-hr reporting)",
          min_len=7),

        P("PIN_CODE",
          re.compile(r'(?<!\d)([1-9]\d{5})(?!\d)'),
          "Indian PIN Code (6-digit)", "LOW",
          "DPDP Act 2023 (location data when linked to address)",
          min_len=6),
    ]

REGEX_PATTERNS = _build_regex_patterns()


# ===================================================================
#  COLUMN HEADER PII KEYWORDS (expanded)
# ===================================================================

COLUMN_PII_KEYWORDS: dict[str, tuple[str, str, str, str]] = {
    # Person Identity
    "employee name":     ("PERSON_NAME",     "HIGH",   "Employee Name",            "DPDP Act 2023, IT Rules 2011, DPDP Rules 2025"),
    "emp name":          ("PERSON_NAME",     "HIGH",   "Employee Name",            "DPDP Act 2023, DPDP Rules 2025"),
    "name":              ("PERSON_NAME",     "HIGH",   "Person Name",              "DPDP Act 2023, DPDP Rules 2025"),
    "first name":        ("PERSON_NAME",     "HIGH",   "First Name",               "DPDP Act 2023, DPDP Rules 2025"),
    "last name":         ("PERSON_NAME",     "HIGH",   "Last Name",                "DPDP Act 2023, DPDP Rules 2025"),
    "father":            ("PERSON_NAME",     "HIGH",   "Father's Name",            "DPDP Act 2023, DPDP Rules 2025"),
    "mother":            ("PERSON_NAME",     "HIGH",   "Mother's Name",            "DPDP Act 2023, DPDP Rules 2025"),
    "reporting manager": ("PERSON_NAME",     "MEDIUM", "Reporting Manager Name",   "DPDP Act 2023"),
    "manager name":      ("PERSON_NAME",     "MEDIUM", "Manager Name",             "DPDP Act 2023"),
    "nominee":           ("PERSON_NAME",     "HIGH",   "Nominee Name",             "DPDP Act 2023, PFRDA Guidelines"),
    "beneficiary":       ("PERSON_NAME",     "HIGH",   "Beneficiary Name",         "DPDP Act 2023, RBI Directions"),
    "guardian":          ("PERSON_NAME",     "HIGH",   "Guardian Name",            "DPDP Act 2023, DPDP Rules 2025"),
    # Dates
    "dob":               ("DATE_OF_BIRTH",   "HIGH",   "Date of Birth",            "DPDP Act 2023, IT Rules 2011, DPDP Rules 2025"),
    "date of birth":     ("DATE_OF_BIRTH",   "HIGH",   "Date of Birth",            "DPDP Act 2023, IT Rules 2011, DPDP Rules 2025"),
    "birth date":        ("DATE_OF_BIRTH",   "HIGH",   "Date of Birth",            "DPDP Act 2023, DPDP Rules 2025"),
    "doj":               ("DATE_EMPLOYMENT", "MEDIUM", "Date of Joining",          "DPDP Act 2023"),
    "date of joining":   ("DATE_EMPLOYMENT", "MEDIUM", "Date of Joining",          "DPDP Act 2023"),
    "doe":               ("DATE_EMPLOYMENT", "MEDIUM", "Date of Exit/Expiry",      "DPDP Act 2023"),
    "doc":               ("DATE_EMPLOYMENT", "MEDIUM", "Date of Confirmation",     "DPDP Act 2023"),
    # Employee
    "employee code":     ("EMPLOYEE_ID",     "MEDIUM", "Employee Code / ID",       "DPDP Act 2023"),
    "emp code":          ("EMPLOYEE_ID",     "MEDIUM", "Employee Code / ID",       "DPDP Act 2023"),
    "employee id":       ("EMPLOYEE_ID",     "MEDIUM", "Employee ID",              "DPDP Act 2023"),
    # Contact
    "email":             ("EMAIL",           "HIGH",   "Email Address",            "DPDP Act 2023, IT Rules 2011, DPDP Rules 2025"),
    "e-mail":            ("EMAIL",           "HIGH",   "Email Address",            "DPDP Act 2023, IT Rules 2011"),
    "phone":             ("PHONE_IN",        "HIGH",   "Phone Number",             "DPDP Act 2023, TRAI Regulations"),
    "mobile":            ("PHONE_IN",        "HIGH",   "Mobile Number",            "DPDP Act 2023, TRAI Regulations"),
    "contact":           ("PHONE_IN",        "HIGH",   "Contact Number",           "DPDP Act 2023, TRAI Regulations"),
    # ID Documents
    "aadhaar":           ("AADHAAR",         "HIGH",   "Aadhaar Number",           "DPDP Act 2023, Aadhaar Act 2016, DPDP Rules 2025"),
    "aadhar":            ("AADHAAR",         "HIGH",   "Aadhaar Number",           "DPDP Act 2023, Aadhaar Act 2016"),
    "uid":               ("AADHAAR",         "HIGH",   "UID / Aadhaar",            "DPDP Act 2023, Aadhaar Act 2016"),
    "pan":               ("PAN_CARD",        "HIGH",   "PAN Card",                 "DPDP Act 2023, Income Tax Act"),
    "pan no":            ("PAN_CARD",        "HIGH",   "PAN Number",               "DPDP Act 2023, Income Tax Act"),
    "passport":          ("PASSPORT_IN",     "HIGH",   "Passport Number",          "DPDP Act 2023, Passports Act 1967"),
    "voter id":          ("VOTER_ID",        "HIGH",   "Voter ID",                 "DPDP Act 2023, RPA 1950"),
    "driving license":   ("DRIVING_LICENSE", "HIGH",   "Driving License",          "DPDP Act 2023, MV Act 1988"),
    "dl no":             ("DRIVING_LICENSE", "HIGH",   "Driving License",          "DPDP Act 2023, MV Act 1988"),
    # Financial
    "salary":            ("FINANCIAL",       "HIGH",   "Salary / Compensation",    "DPDP Act 2023, Payment of Wages Act"),
    "ctc":               ("FINANCIAL",       "HIGH",   "CTC / Compensation",       "DPDP Act 2023, Payment of Wages Act"),
    "gross":             ("FINANCIAL",       "HIGH",   "Gross Salary",             "DPDP Act 2023, Payment of Wages Act"),
    "net pay":           ("FINANCIAL",       "HIGH",   "Net Pay",                  "DPDP Act 2023, Payment of Wages Act"),
    "basic pay":         ("FINANCIAL",       "HIGH",   "Basic Pay",                "DPDP Act 2023, Payment of Wages Act"),
    "bank account":      ("BANK_ACCOUNT",    "HIGH",   "Bank Account Number",      "RBI Master Directions, DPDP Act 2023"),
    "account no":        ("BANK_ACCOUNT",    "HIGH",   "Account Number",           "RBI Master Directions, DPDP Act 2023"),
    "ifsc":              ("IFSC_CODE",       "MEDIUM", "IFSC Code",                "RBI Master Directions"),
    "gstin":             ("GSTIN",           "MEDIUM", "GSTIN",                    "GST Act 2017"),
    "gst no":            ("GSTIN",           "MEDIUM", "GST Number",               "GST Act 2017"),
    # Insurance & Pension
    "policy no":         ("INSURANCE",       "HIGH",   "Insurance Policy Number",  "IRDAI Cybersecurity Guidelines 2023, DPDP Act 2023"),
    "policy number":     ("INSURANCE",       "HIGH",   "Insurance Policy Number",  "IRDAI Cybersecurity Guidelines 2023, DPDP Act 2023"),
    "pran":              ("PENSION",         "HIGH",   "PRAN (Pension Account)",   "PFRDA Subscriber Data Guidelines, DPDP Act 2023"),
    "pension":           ("PENSION",         "HIGH",   "Pension Data",             "PFRDA Subscriber Data Guidelines, DPDP Act 2023"),
    "provident fund":    ("PENSION",         "HIGH",   "Provident Fund Data",      "PFRDA Guidelines, DPDP Act 2023"),
    "pf no":             ("PENSION",         "HIGH",   "PF Number",               "PFRDA Guidelines, DPDP Act 2023"),
    "uan":               ("PENSION",         "HIGH",   "UAN (Universal Account)",  "EPFO, DPDP Act 2023"),
    # Health
    "blood group":       ("HEALTH_DATA",     "HIGH",   "Blood Group",              "DPDP Act 2023 Sec 2(t), IT Rules 2011 Sec 3"),
    "medical":           ("HEALTH_DATA",     "HIGH",   "Medical Information",       "DPDP Act 2023, IT Rules 2011 Sec 3"),
    "health":            ("HEALTH_DATA",     "HIGH",   "Health Data",              "DPDP Act 2023, IT Rules 2011 Sec 3"),
    "disability":        ("HEALTH_DATA",     "HIGH",   "Disability Information",    "DPDP Act 2023, RPwD Act 2016"),
    # Location & Address
    "address":           ("ADDRESS",         "MEDIUM", "Address / Location",        "DPDP Act 2023, DPDP Rules 2025"),
    "residential":       ("ADDRESS",         "MEDIUM", "Residential Address",       "DPDP Act 2023"),
    "permanent address": ("ADDRESS",         "MEDIUM", "Permanent Address",         "DPDP Act 2023"),
    "pin code":          ("PIN_CODE",        "LOW",    "PIN Code",                  "DPDP Act 2023"),
    "pincode":           ("PIN_CODE",        "LOW",    "PIN Code",                  "DPDP Act 2023"),
    # Biometric
    "biometric":         ("BIOMETRIC",       "HIGH",   "Biometric Data",            "DPDP Act 2023, Aadhaar Act 2016 Sec 29"),
    "fingerprint":       ("BIOMETRIC",       "HIGH",   "Fingerprint Data",          "DPDP Act 2023, Aadhaar Act 2016 Sec 29"),
    # Other
    "payroll":           ("FINANCIAL",       "HIGH",   "Payroll Data",              "DPDP Act 2023, Payment of Wages Act"),
    "cost center":       ("ORG_DATA",        "LOW",    "Cost Center / Org Data",    "Internal Policy"),
    "department":        ("ORG_DATA",        "LOW",    "Department",                "Internal Policy"),
    "designation":       ("ORG_DATA",        "LOW",    "Designation / Role",        "Internal Policy"),
    # Securities / SEBI
    "demat":             ("DEMAT_ACCOUNT",   "HIGH",   "Demat Account Number",      "SEBI CSCRF 2023, DPDP Act 2023"),
    "dp id":             ("DEMAT_ACCOUNT",   "HIGH",   "DP ID",                     "SEBI CSCRF 2023, DPDP Act 2023"),
    "client id":         ("CLIENT_ID",       "MEDIUM", "Client ID",                 "SEBI CSCRF 2023, DPDP Act 2023"),
    "folio":             ("FOLIO_NUMBER",     "MEDIUM", "Folio Number",              "SEBI CSCRF 2023, DPDP Act 2023"),
}


# ===================================================================
#  VALUE MASKING
# ===================================================================

def _mask_value(tag: str, value: str) -> str:
    v = str(value).strip()
    if not v:
        return "***"
    if tag == "AADHAAR":
        d = re.sub(r'\D', '', v)
        return f"XXXX XXXX {d[-4:]}" if len(d) >= 4 else "XXXX XXXX XXXX"
    if tag == "PAN_CARD":
        return f"{v[:2]}XXX{v[5:8]}XX{v[-1]}" if len(v) == 10 else "XXXXXXXXXX"
    if tag == "EMAIL":
        parts = v.split("@")
        if len(parts) == 2:
            u = parts[0]
            return f"{u[0]}***{u[-1] if len(u) > 1 else ''}@{parts[1]}"
        return "***@***"
    if tag == "PHONE_IN":
        d = re.sub(r'\D', '', v)
        return f"XXXXXX{d[-4:]}" if len(d) >= 4 else "XXXXXXXXXX"
    if tag == "CREDIT_CARD":
        d = re.sub(r'\D', '', v)
        return f"XXXX-XXXX-XXXX-{d[-4:]}" if len(d) >= 4 else "XXXX-XXXX-XXXX-XXXX"
    if tag == "IP_ADDRESS":
        parts = v.split(".")
        return f"{parts[0]}.{parts[1]}.XXX.XXX" if len(parts) == 4 else "XXX.XXX.XXX.XXX"
    if tag in ("PASSPORT_IN", "VOTER_ID", "DRIVING_LICENSE"):
        return v[:2] + "X" * (len(v) - 3) + v[-1] if len(v) > 3 else "X" * len(v)
    if tag == "UPI_ID":
        parts = v.split("@")
        return f"{parts[0][:2]}***@{parts[1]}" if len(parts) == 2 else "***@***"
    if tag == "PERSON_NAME":
        words = v.split()
        masked = []
        for w in words:
            if len(w) > 2:
                masked.append(w[0] + "*" * (len(w) - 2) + w[-1])
            else:
                masked.append(w[0] + "*")
        return " ".join(masked)
    # Default
    if len(v) > 4:
        return v[:2] + "X" * (len(v) - 4) + v[-2:]
    return "X" * len(v)


# ===================================================================
#  FILE PARSERS
# ===================================================================

def _compute_sha256(filepath: str) -> str:
    h = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return "HASH_ERROR"


def _count_docx_images(filepath: str) -> int:
    count = 0
    try:
        with zipfile.ZipFile(filepath) as z:
            count = sum(1 for n in z.namelist() if 'media/' in n and
                        any(n.lower().endswith(e) for e in
                            ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')))
    except Exception:
        pass
    return count


def _count_pdf_images(filepath: str) -> int:
    count = 0
    if not HAS_PYPDF:
        return 0
    try:
        reader = pypdf.PdfReader(filepath)
        for page in reader.pages:
            resources = page.get("/Resources")
            if resources:
                xobjects = resources.get("/XObject")
                if xobjects:
                    try:
                        xobj = xobjects.get_object()
                        for name in xobj:
                            obj = xobj[name].get_object()
                            if obj.get("/Subtype") == "/Image":
                                count += 1
                    except Exception:
                        pass
    except Exception:
        pass
    return count


def _ocr_image(filepath: str) -> str | None:
    if not HAS_PIL:
        return None
    try:
        img = Image.open(filepath)
        if HAS_TESSERACT:
            text = pytesseract.image_to_string(img, lang="eng")
            return _sanitize(text) if text.strip() else None
        return None
    except Exception:
        return None


def extract_text(filepath: str) -> tuple[str | None, str, str | None, int]:
    """Returns (text, file_type, error, embedded_image_count)."""
    ext = pathlib.Path(filepath).suffix.lower()
    img_count = 0

    try:
        if ext in IMAGE_EXTENSIONS:
            ocr_text = _ocr_image(filepath)
            if ocr_text:
                return ocr_text, "IMAGE", None, 1
            note = "Image file (OCR not available)" if not HAS_TESSERACT else "Image (no text found)"
            return None, "IMAGE", note, 1

        if ext in {".txt", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml",
                   ".log", ".sql", ".env", ".ini", ".conf", ".cfg", ".md",
                   ".html", ".htm", ".rtf"}:
            for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
                try:
                    with open(filepath, "r", encoding=enc, errors="replace") as f:
                        text = f.read()
                    return _sanitize(text), ext.lstrip(".").upper(), None, 0
                except (UnicodeDecodeError, UnicodeError):
                    continue
            return None, ext.lstrip(".").upper(), "Encoding error", 0

        if ext == ".pdf":
            if not HAS_PYPDF:
                return None, "PDF", "pypdf not installed", 0
            reader = pypdf.PdfReader(filepath)
            pages = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            img_count = _count_pdf_images(filepath)
            return _sanitize("\n".join(pages)), "PDF", None, img_count

        if ext == ".docx":
            if not HAS_DOCX:
                return None, "DOCX", "python-docx not installed", 0
            doc = DocxDocument(filepath)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
            img_count = _count_docx_images(filepath)
            return _sanitize("\n".join(parts)), "DOCX", None, img_count

        if ext in (".xlsx", ".xls"):
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            col_findings = []
            for ws in wb.worksheets:
                headers = []
                col_samples = defaultdict(list)
                total_rows = ws.max_row or 1

                for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                    row_vals = [str(c) if c is not None else "" for c in row]
                    parts.append(" | ".join(row_vals))

                    if row_idx == 1:
                        headers = [v.strip() for v in row_vals]
                    elif row_idx <= 22:
                        for col_idx, val in enumerate(row_vals):
                            if val.strip() and len(col_samples[col_idx]) < 3:
                                col_samples[col_idx].append(val)

                for col_idx, header in enumerate(headers):
                    hl = header.lower().strip()
                    if not hl:
                        continue
                    for keyword, (tag, sens, desc, reg) in COLUMN_PII_KEYWORDS.items():
                        if keyword in hl:
                            samples = col_samples.get(col_idx, [])
                            if samples:
                                sample_masked = ", ".join([_mask_value(tag, s) for s in samples])
                                col_findings.append({
                                    "tag": tag,
                                    "description": desc,
                                    "sensitivity": sens,
                                    "regulation": reg,
                                    "raw_value": f"Column: '{header}'",
                                    "masked_value": f"Column: '{header}' ({total_rows - 1} rows)",
                                    "line_number": 0,
                                    "context": f"Sheet: {ws.title} | Header: {header} | Samples: {sample_masked}",
                                    "confidence": 95,
                                    "detection_method": "Column Header Analysis",
                                })
                            break

            wb.close()
            return _sanitize("\n".join(parts)), "XLSX", None, 0, col_findings

        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                return _sanitize("\n".join(texts)), "PPTX", None, 0
            except ImportError:
                return None, "PPTX", "python-pptx not installed", 0

        return None, ext.lstrip(".").upper(), "Unsupported type", 0

    except PermissionError:
        return None, ext.lstrip(".").upper(), "Permission denied", 0
    except Exception as e:
        return None, ext.lstrip(".").upper(), str(e)[:100], 0


# ===================================================================
#  LAYER 1 — REGEX SCANNING
# ===================================================================

def scan_regex(text: str) -> list[dict]:
    findings = []
    if not text:
        return findings

    # Pre-compute line start offsets for fast binary-search line number mapping
    line_starts = [0] + [m.start() + 1 for m in re.finditer(r'\n', text)]

    for pat in REGEX_PATTERNS:
        for match in pat.regex.finditer(text):
            raw = match.group().strip()
            if len(raw) < pat.min_len:
                continue

            checksum_valid = None
            if pat.validator:
                checksum_valid = pat.validator(match)
                if not checksum_valid:
                    continue

            start_pos = match.start()
            end_pos = match.end()

            line_num = bisect.bisect_right(line_starts, start_pos)
            line_start = line_starts[line_num - 1]
            line_end = line_starts[line_num] - 1 if line_num < len(line_starts) else len(text)
            line = text[line_start:line_end]

            rel_start = start_pos - line_start
            rel_end = end_pos - line_start
            ctx_start = max(0, rel_start - 40)
            ctx_end = min(len(line), rel_end + 40)
            ctx = line[ctx_start:ctx_end].strip()

            masked = _mask_value(pat.tag, raw)
            ctx_masked = _sanitize(ctx.replace(raw, masked))

            confidence = 85
            if checksum_valid is True:
                confidence = 99
            if pat.tag == "IP_ADDRESS":
                confidence = 70
                if raw.startswith("127.") or raw.startswith("0."):
                    confidence = 40
            if pat.tag == "PIN_CODE":
                confidence = 50

            findings.append({
                "tag": pat.tag,
                "description": pat.description,
                "sensitivity": pat.sensitivity,
                "regulation": pat.regulation,
                "raw_value": raw,
                "masked_value": masked,
                "line_number": line_num,
                "context": ctx_masked,
                "confidence": confidence,
                "detection_method": "Regex" + (" + Checksum" if checksum_valid else ""),
            })

    return findings


# ===================================================================
#  LAYER 2 — spaCy NER SCANNING
# ===================================================================

def scan_ner(text: str) -> list[dict]:
    nlp, model_name = get_nlp()
    if not HAS_SPACY or nlp is None:
        return []

    findings = []
    max_chunk = 100000
    chunks = [text[i:i + max_chunk] for i in range(0, len(text), max_chunk)]
    chunk_offsets = [i * max_chunk for i in range(len(chunks))]

    skip_words = {
        "dear", "sir", "madam", "mr", "mrs", "ms", "regards",
        "thanks", "hi", "hello", "subject", "from", "to", "cc",
        "sent", "date", "re", "fw", "fwd", "working", "ok",
        "testing", "test", "related", "close", "arrange",
        "none", "null", "n/a", "na", "true", "false",
    }

    line_starts = [0] + [m.start() + 1 for m in re.finditer(r'\n', text)]

    for doc, chunk_offset in zip(nlp.pipe(chunks, batch_size=8), chunk_offsets):
        chunk_text = doc.text
        for ent in doc.ents:
            tag = None
            sensitivity = "MEDIUM"
            description = ""
            regulation = "DPDP Act 2023, DPDP Rules 2025"

            if ent.label_ == "PERSON":
                if len(ent.text.strip()) < 3:
                    continue
                lower = ent.text.strip().lower()
                if lower in skip_words:
                    continue
                tag = "PERSON_NAME"
                sensitivity = "HIGH"
                description = "Person Name (NER Detected)"
                regulation = "DPDP Act 2023 Sec 2(t), IT Rules 2011 Sec 3, DPDP Rules 2025"

            elif ent.label_ == "ORG":
                if len(ent.text.strip()) < 3:
                    continue
                tag = "ORGANISATION"
                sensitivity = "LOW"
                description = "Organisation Name (NER Detected)"
                regulation = "DPDP Act 2023, Companies Act 2013"

            elif ent.label_ == "GPE":
                if len(ent.text.strip()) < 3:
                    continue
                tag = "LOCATION"
                sensitivity = "LOW"
                description = "Location / Address (NER Detected)"
                regulation = "DPDP Act 2023 Sec 2(t)"

            if tag:
                masked = _mask_value(tag, ent.text.strip())
                abs_start = chunk_offset + ent.start_char
                abs_end = chunk_offset + ent.end_char
                line_num = bisect.bisect_right(line_starts, abs_start)

                ctx_start = max(0, ent.start_char - 40)
                ctx_end = min(len(chunk_text), ent.end_char + 40)
                ctx = _sanitize(chunk_text[ctx_start:ctx_end].replace(ent.text, masked))

                confidence = 75
                if tag == "PERSON_NAME" and len(ent.text.split()) >= 2:
                    confidence = 85
                if tag in ("ORGANISATION", "LOCATION"):
                    confidence = 65
                if model_name == "en_core_web_lg":
                    confidence = min(99, confidence + 5)

                findings.append({
                    "tag": tag,
                    "description": description,
                    "sensitivity": sensitivity,
                    "regulation": regulation,
                    "raw_value": ent.text.strip(),
                    "masked_value": masked,
                    "line_number": line_num,
                    "context": ctx,
                    "confidence": confidence,
                    "detection_method": f"spaCy NER ({model_name})",
                })

    return findings


# ===================================================================
#  LAYER 3 — COLUMN HEADER ANALYSIS
# ===================================================================

def scan_columns_xlsx(filepath: str) -> list[dict]:
    findings = []
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    except Exception:
        return findings

    for ws in wb.worksheets:
        headers = []
        for row in ws.iter_rows(min_row=1, max_row=1, values_only=True):
            headers = [str(h).strip() if h else "" for h in row]
            break

        for col_idx, header in enumerate(headers):
            hl = header.lower().strip()
            for keyword, (tag, sens, desc, reg) in COLUMN_PII_KEYWORDS.items():
                if keyword in hl:
                    data_count = 0
                    samples = []
                    for row in ws.iter_rows(min_row=2,
                                            max_row=min((ws.max_row or 2), 22),
                                            min_col=col_idx + 1,
                                            max_col=col_idx + 1,
                                            values_only=True):
                        val = row[0]
                        if val is not None and str(val).strip():
                            data_count += 1
                            if len(samples) < 3:
                                samples.append(str(val))

                    if data_count > 0:
                        total = (ws.max_row or 1) - 1
                        sample_masked = ", ".join(
                            [_mask_value(tag, s) for s in samples])
                        findings.append({
                            "tag": tag,
                            "description": desc,
                            "sensitivity": sens,
                            "regulation": reg,
                            "raw_value": f"Column: '{header}'",
                            "masked_value": f"Column: '{header}' ({total} rows)",
                            "line_number": 0,
                            "context": f"Sheet: {ws.title} | Header: {header} | "
                                       f"Samples: {sample_masked}",
                            "confidence": 95,
                            "detection_method": "Column Header Analysis",
                        })
                    break

    wb.close()
    return findings


def scan_columns_csv(filepath: str) -> list[dict]:
    findings = []
    headers = []
    sample_rows = []
    try:
        for enc in ("utf-8", "utf-8-sig", "latin-1"):
            try:
                with open(filepath, "r", encoding=enc, errors="replace") as f:
                    reader = csv.reader(f)
                    headers = next(reader, [])
                    if not headers:
                        return findings
                    for i, row in enumerate(reader):
                        sample_rows.append(row)
                        if i >= 20:
                            break
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
    except Exception:
        return findings

    for col_idx, header in enumerate(headers):
        hl = header.lower().strip()
        for keyword, (tag, sens, desc, reg) in COLUMN_PII_KEYWORDS.items():
            if keyword in hl:
                samples = []
                for row in sample_rows:
                    if col_idx < len(row) and row[col_idx].strip():
                        samples.append(row[col_idx])
                        if len(samples) >= 3:
                            break
                if samples:
                    sample_masked = ", ".join(
                        [_mask_value(tag, s) for s in samples])
                    findings.append({
                        "tag": tag,
                        "description": desc,
                        "sensitivity": sens,
                        "regulation": reg,
                        "raw_value": f"Column: '{header}'",
                        "masked_value": f"Column: '{header}'",
                        "line_number": 0,
                        "context": f"CSV Header: {header} | Samples: {sample_masked}",
                        "confidence": 95,
                        "detection_method": "Column Header Analysis",
                    })
                break

    return findings


# ===================================================================
#  DEDUPLICATION
# ===================================================================

def deduplicate_findings(findings: list[dict]) -> list[dict]:
    seen: dict[str, dict] = {}
    for f in findings:
        key = f"{f['tag']}|{f.get('raw_value', '')}|{f['line_number']}"
        if key in seen:
            existing = seen[key]
            if f["detection_method"] != existing["detection_method"]:
                existing["confidence"] = min(99, existing["confidence"] + 15)
                existing["detection_method"] += f" + {f['detection_method']}"
            elif f["confidence"] > existing["confidence"]:
                seen[key] = f
        else:
            seen[key] = f
    return list(seen.values())


# ===================================================================
#  FILE SCANNER — ORCHESTRATOR
# ===================================================================

def scan_single_file(fp: pathlib.Path, target: pathlib.Path) -> tuple[list[dict], dict]:
    """Scans a single file using all detection layers. Returns (file_findings, audit_record)."""
    rel = str(fp.relative_to(target)) if fp.is_relative_to(target) else str(fp)
    try:
        stat = fp.stat()
        fsize = stat.st_size
        fmod = datetime.datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
    except Exception:
        fsize = 0
        fmod = "N/A"

    file_hash = _compute_sha256(str(fp))
    extracted = extract_text(str(fp))
    text, ftype, error, img_count = extracted[0], extracted[1], extracted[2], extracted[3]
    col_findings = extracted[4] if len(extracted) > 4 else []

    if error:
        audit_record = {
            "file_name": fp.name, "file_path": str(fp),
            "relative_path": rel, "file_type": ftype,
            "file_size": fsize, "last_modified": fmod,
            "sha256": file_hash, "status": "SKIPPED",
            "reason": error, "pii_tags": "", "pii_count": 0,
            "risk_score": "N/A", "embedded_images": img_count,
        }
        return [], audit_record

    # ------ Detection Pipeline ------
    file_findings = []

    # Layer 1: Regex
    if text:
        file_findings.extend(scan_regex(text))

    # Layer 2: spaCy NER (only for unstructured text)
    is_sheet = fp.suffix.lower() in (".xlsx", ".xls", ".csv", ".tsv")
    if text and not is_sheet:
        file_findings.extend(scan_ner(text))

    # Layer 3: Column header analysis
    if col_findings:
        file_findings.extend(col_findings)
    elif fp.suffix.lower() in (".csv", ".tsv"):
        file_findings.extend(scan_columns_csv(str(fp)))

    # Embedded images flag
    if img_count > 0:
        file_findings.append({
            "tag": "EMBEDDED_IMAGES",
            "description": (f"{img_count} embedded image(s) — may contain "
                            "PII (screenshots, IDs, scanned documents)"),
            "sensitivity": "MEDIUM",
            "regulation": "DPDP Act 2023, IT Rules 2011, DPDP Rules 2025",
            "raw_value": f"{img_count} image(s)",
            "masked_value": f"{img_count} embedded image(s)",
            "line_number": 0,
            "context": (f"File contains {img_count} embedded image(s). "
                        "Manual review recommended."),
            "confidence": 60,
            "detection_method": "File Structure Analysis",
        })

    # Deduplicate
    file_findings = deduplicate_findings(file_findings)
    pii_count = len(file_findings)

    # Attach file metadata
    for f in file_findings:
        f["file_name"] = fp.name
        f["file_path"] = str(fp)
        f["relative_path"] = rel
        f["file_type"] = ftype
        f["file_size"] = fsize
        f["last_modified"] = fmod
        f["sha256"] = file_hash

    if pii_count == 0:
        risk_score = "CLEAN"
    elif any(f["sensitivity"] == "HIGH" for f in file_findings):
        risk_score = "HIGH"
    elif any(f["sensitivity"] == "MEDIUM" for f in file_findings):
        risk_score = "MEDIUM"
    else:
        risk_score = "LOW"

    tags_found = sorted(set(f["tag"] for f in file_findings))
    status = "PII_DETECTED" if pii_count > 0 else "CLEAN"

    audit_record = {
        "file_name": fp.name, "file_path": str(fp),
        "relative_path": rel, "file_type": ftype,
        "file_size": fsize, "last_modified": fmod,
        "sha256": file_hash, "status": status,
        "reason": (f"{pii_count} finding(s)" if pii_count > 0
                   else "No PII detected"),
        "pii_tags": ", ".join(tags_found),
        "pii_count": pii_count, "risk_score": risk_score,
        "embedded_images": img_count,
    }

    return file_findings, audit_record


def scan_folder(target_path: str) -> tuple[list[dict], list[dict]]:
    all_findings: list[dict] = []
    file_audit: list[dict] = []
    target = pathlib.Path(target_path)

    if not target.exists():
        print(f"  [-] Target path does not exist: {target}")
        return all_findings, file_audit

    files_to_scan = []
    if target.is_file():
        files_to_scan = [target]
    else:
        for root, dirs, filenames in os.walk(target):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for fn in filenames:
                if fn.startswith("~$"):
                    continue
                fp = pathlib.Path(root) / fn
                if fp.suffix.lower() in SCAN_EXTENSIONS:
                    files_to_scan.append(fp)

    total = len(files_to_scan)
    print(f"  [+] {total} file(s) to scan in: {target}")
    print()

    for idx, fp in enumerate(files_to_scan, 1):
        file_findings, audit_record = scan_single_file(fp, target)
        all_findings.extend(file_findings)
        file_audit.append(audit_record)

        rel = audit_record["relative_path"]
        pii_count = audit_record["pii_count"]
        tags_found = audit_record["pii_tags"]
        img_count = audit_record["embedded_images"]
        status = audit_record["status"]

        if status == "SKIPPED":
            print(f"  [{idx}/{total}] [SKIP] {rel} — {audit_record['reason']}")
        elif pii_count > 0:
            img_note = f" | Images: {img_count}" if img_count > 0 else ""
            print(f"  [{idx}/{total}] [!] PII DETECTED: {rel} ({pii_count} findings) | Tags: {tags_found}{img_note}")
        else:
            print(f"  [{idx}/{total}] [OK] Clean: {rel}")

    return all_findings, file_audit

    return all_findings, file_audit


# ===================================================================
#  EXCEL REPORT GENERATOR — 4 PROFESSIONAL TABS
# ===================================================================

def generate_report(findings: list[dict], file_audit: list[dict],
                    output_path: str, target_path: str) -> str:
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    # -- Styles --
    NAVY = PatternFill("solid", fgColor="1B2A4A")
    DARK = PatternFill("solid", fgColor="2C3E50")
    GREY = PatternFill("solid", fgColor="F8F9FA")

    H_FONT = Font("Calibri", 11, bold=True, color="FFFFFF")
    T_FONT = Font("Calibri", 16, bold=True, color="1B2A4A")
    S_FONT = Font("Calibri", 10, italic=True, color="6C757D")
    L_FONT = Font("Calibri", 9, bold=True, color="6C757D")
    V_FONT = Font("Calibri", 22, bold=True, color="1B2A4A")
    SEC_F = Font("Calibri", 12, bold=True, color="1B2A4A")
    D_FONT = Font("Calibri", 10, color="212529")
    B_FONT = Font("Calibri", 10, bold=True, color="212529")

    THIN = Side(border_style="thin", color="DEE2E6")
    BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

    RISK_STYLE = {
        "HIGH":   (PatternFill("solid", fgColor="F8D7DA"),
                   Font("Calibri", 10, bold=True, color="842029")),
        "MEDIUM": (PatternFill("solid", fgColor="FFF3CD"),
                   Font("Calibri", 10, bold=True, color="664D03")),
        "LOW":    (PatternFill("solid", fgColor="D1E7DD"),
                   Font("Calibri", 10, bold=True, color="0F5132")),
    }
    STATUS_FILL = {
        "PII_DETECTED": PatternFill("solid", fgColor="F8D7DA"),
        "CLEAN":        PatternFill("solid", fgColor="D1E7DD"),
        "SKIPPED":      PatternFill("solid", fgColor="E2E3E5"),
        "HIGH":         PatternFill("solid", fgColor="F8D7DA"),
        "MEDIUM":       PatternFill("solid", fgColor="FFF3CD"),
        "LOW":          PatternFill("solid", fgColor="D1E7DD"),
        "N/A":          PatternFill("solid", fgColor="E2E3E5"),
    }

    def _hdr(ws, row, headers, fill=NAVY):
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=row, column=c, value=h)
            cell.fill = fill
            cell.font = H_FONT
            cell.alignment = Alignment(horizontal="center", vertical="center",
                                       wrap_text=True)
            cell.border = BORDER

    def _dcell(cell, wrap=False):
        cell.font = D_FONT
        cell.border = BORDER
        cell.alignment = Alignment(vertical="center", wrap_text=wrap)

    def _autofit(ws, mx=55):
        for col in ws.columns:
            ml = max((len(str(c.value or "")) for c in col), default=10)
            ws.column_dimensions[get_column_letter(
                col[0].column)].width = min(max(ml + 3, 12), mx)

    # -- Metrics --
    total_findings = len(findings)
    total_files = len(file_audit)
    files_pii = sum(1 for f in file_audit if f["status"] == "PII_DETECTED")
    files_clean = sum(1 for f in file_audit if f["status"] == "CLEAN")
    files_skip = sum(1 for f in file_audit if f["status"] == "SKIPPED")
    high_c = sum(1 for f in findings if f["sensitivity"] == "HIGH")
    med_c = sum(1 for f in findings if f["sensitivity"] == "MEDIUM")
    low_c = sum(1 for f in findings if f["sensitivity"] == "LOW")
    avg_conf = round(
        sum(f["confidence"] for f in findings) / max(total_findings, 1), 1)
    scan_ts = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    tag_freq = defaultdict(int)
    for f in findings:
        tag_freq[f["tag"]] += 1
    sorted_tags = sorted(tag_freq.items(), key=lambda x: x[1], reverse=True)

    method_freq = defaultdict(int)
    for f in findings:
        for m in f["detection_method"].split(" + "):
            method_freq[m.strip()] += 1

    # ================================================================
    #  TAB 1 — DASHBOARD
    # ================================================================
    ws1 = wb.create_sheet("Dashboard")

    ws1.merge_cells("A1:J1")
    ws1["A1"] = "PII DATA DISCOVERY — AUDIT REPORT"
    ws1["A1"].font = T_FONT

    ws1.merge_cells("A2:J2")
    ws1["A2"] = (f"Scan: {scan_ts}  |  Target: {target_path}  |  "
                 f"v{TOOL_VERSION}  |  NER: {_NLP_MODEL_NAME}")
    ws1["A2"].font = S_FONT

    # KPI Row 1
    kpis_1 = [
        ("FILES SCANNED", total_files,   "A", "B", "E8EAF6"),
        ("FILES AT RISK", files_pii,     "C", "D", "F8D7DA"),
        ("TOTAL FINDINGS", total_findings, "E", "F", "FFF3CD"),
        ("HIGH RISK",     high_c,        "G", "H", "F8D7DA"),
        ("AVG CONFIDENCE", f"{avg_conf}%", "I", "J", "D1ECF1"),
    ]
    for label, val, c1, c2, color in kpis_1:
        fill = PatternFill("solid", fgColor=color)
        ws1.merge_cells(f"{c1}4:{c2}4")
        cl = ws1[f"{c1}4"]
        cl.value = label
        cl.font = L_FONT
        cl.alignment = Alignment(horizontal="center", vertical="center")
        ws1.merge_cells(f"{c1}5:{c2}5")
        cv = ws1[f"{c1}5"]
        cv.value = val
        cv.font = V_FONT
        cv.alignment = Alignment(horizontal="center", vertical="center")
        for r in (4, 5):
            for letter in (c1, c2):
                ws1[f"{letter}{r}"].fill = fill
                ws1[f"{letter}{r}"].border = BORDER

    # KPI Row 2
    kpis_2 = [
        ("MEDIUM RISK",    med_c,       "A", "B", "FFF3CD"),
        ("LOW RISK",       low_c,       "C", "D", "D1E7DD"),
        ("CLEAN FILES",    files_clean, "E", "F", "D1E7DD"),
        ("SKIPPED FILES",  files_skip,  "G", "H", "E2E3E5"),
        ("PII CATEGORIES", len(tag_freq), "I", "J", "E8EAF6"),
    ]
    for label, val, c1, c2, color in kpis_2:
        fill = PatternFill("solid", fgColor=color)
        ws1.merge_cells(f"{c1}7:{c2}7")
        cl = ws1[f"{c1}7"]
        cl.value = label
        cl.font = L_FONT
        cl.alignment = Alignment(horizontal="center", vertical="center")
        ws1.merge_cells(f"{c1}8:{c2}8")
        cv = ws1[f"{c1}8"]
        cv.value = val
        cv.font = V_FONT
        cv.alignment = Alignment(horizontal="center", vertical="center")
        for r in (7, 8):
            for letter in (c1, c2):
                ws1[f"{letter}{r}"].fill = fill
                ws1[f"{letter}{r}"].border = BORDER

    # Section: PII Categories Found
    r = 10
    ws1.cell(row=r, column=1, value="PII Categories Found").font = SEC_F
    r += 1
    _hdr(ws1, r, ["PII Category", "Count", "Risk Level",
                   "", "", "", "", "", "", ""], DARK)
    for i, (tag, cnt) in enumerate(sorted_tags, r + 1):
        ws1.cell(row=i, column=1, value=tag).border = BORDER
        ws1.cell(row=i, column=2, value=cnt).border = BORDER
        # Determine risk for this tag
        tag_risk = "LOW"
        for f in findings:
            if f["tag"] == tag:
                tag_risk = f["sensitivity"]
                break
        rc = ws1.cell(row=i, column=3, value=tag_risk)
        rc.border = BORDER
        if tag_risk in RISK_STYLE:
            rc.fill, rc.font = RISK_STYLE[tag_risk]
        rc.alignment = Alignment(horizontal="center")

    # Section: File Risk Summary
    r2 = max(r + len(sorted_tags) + 3, 14)
    ws1.cell(row=r2, column=1, value="File Risk Summary").font = SEC_F
    r2 += 1
    _hdr(ws1, r2, ["File Name", "Status", "Risk", "PII Count",
                    "PII Types", "Images", "", "", "", ""], DARK)
    for i, entry in enumerate(file_audit, r2 + 1):
        ws1.cell(row=i, column=1, value=entry["file_name"]).border = BORDER
        st = ws1.cell(row=i, column=2, value=entry["status"])
        st.border = BORDER
        if entry["status"] in STATUS_FILL:
            st.fill = STATUS_FILL[entry["status"]]
        st.alignment = Alignment(horizontal="center")
        rs = ws1.cell(row=i, column=3, value=entry["risk_score"])
        rs.border = BORDER
        if entry["risk_score"] in STATUS_FILL:
            rs.fill = STATUS_FILL[entry["risk_score"]]
        rs.alignment = Alignment(horizontal="center")
        ws1.cell(row=i, column=4, value=entry["pii_count"]).border = BORDER
        ws1.cell(row=i, column=5, value=entry["pii_tags"]).border = BORDER
        imgs = entry.get("embedded_images", 0)
        ws1.cell(row=i, column=6,
                 value=imgs if imgs > 0 else "-").border = BORDER

    # Section: Detection Methods
    r3 = r2 + len(file_audit) + 3
    ws1.cell(row=r3, column=1, value="Detection Methods Used").font = SEC_F
    r3 += 1
    _hdr(ws1, r3, ["Method", "Findings", "", "", "", "", "", "", "", ""], DARK)
    for i, (method, cnt) in enumerate(
            sorted(method_freq.items(), key=lambda x: -x[1]), r3 + 1):
        ws1.cell(row=i, column=1, value=method).border = BORDER
        ws1.cell(row=i, column=2, value=cnt).border = BORDER

    _autofit(ws1)

    # ================================================================
    #  TAB 2 — FINDINGS
    # ================================================================
    ws2 = wb.create_sheet("Findings")

    hdrs2 = [
        "#", "File", "PII Type", "What Was Found",
        "Risk", "Confidence", "How Detected",
        "Where", "Indian Law", "Context"
    ]
    _hdr(ws2, 1, hdrs2)

    for ri, f in enumerate(findings, 2):
        ws2.cell(row=ri, column=1, value=ri - 1)
        ws2.cell(row=ri, column=2, value=f["file_name"])
        ws2.cell(row=ri, column=3, value=f["tag"])
        ws2.cell(row=ri, column=4, value=_sanitize(f["masked_value"]))

        rc = ws2.cell(row=ri, column=5, value=f["sensitivity"])
        if f["sensitivity"] in RISK_STYLE:
            rc.fill, rc.font = RISK_STYLE[f["sensitivity"]]
            rc.alignment = Alignment(horizontal="center")

        ws2.cell(row=ri, column=6, value=f"{f['confidence']}%")
        ws2.cell(row=ri, column=7, value=f["detection_method"])

        loc = (f"Line {f['line_number']}" if f["line_number"] > 0
               else "Column-level")
        ws2.cell(row=ri, column=8, value=loc)
        ws2.cell(row=ri, column=9, value=f["regulation"])
        ws2.cell(row=ri, column=10, value=_sanitize(f.get("context", "")))

        for c in range(1, 11):
            _dcell(ws2.cell(row=ri, column=c), wrap=(c in (4, 9, 10)))
        if ri % 2 == 0:
            for c in range(1, 11):
                if c != 5:
                    ws2.cell(row=ri, column=c).fill = GREY

    if findings:
        ws2.auto_filter.ref = f"A1:J{len(findings) + 1}"
    _autofit(ws2)

    # ================================================================
    #  TAB 3 — FILE AUDIT LOG
    # ================================================================
    ws3 = wb.create_sheet("File Audit Log")

    hdrs3 = [
        "#", "File Name", "Path", "Type", "Size (KB)",
        "Status", "Risk", "PII Types Found", "Findings",
        "Images", "SHA-256", "Last Modified"
    ]
    _hdr(ws3, 1, hdrs3)

    for ri, entry in enumerate(file_audit, 2):
        skb = round(entry["file_size"] / 1024, 1)
        ws3.cell(row=ri, column=1, value=ri - 1)
        ws3.cell(row=ri, column=2, value=entry["file_name"])
        ws3.cell(row=ri, column=3, value=entry["relative_path"])
        ws3.cell(row=ri, column=4, value=entry["file_type"])
        ws3.cell(row=ri, column=5, value=skb)

        st = ws3.cell(row=ri, column=6, value=entry["status"])
        if entry["status"] in STATUS_FILL:
            st.fill = STATUS_FILL[entry["status"]]
        st.alignment = Alignment(horizontal="center")

        rs = ws3.cell(row=ri, column=7, value=entry["risk_score"])
        if entry["risk_score"] in STATUS_FILL:
            rs.fill = STATUS_FILL[entry["risk_score"]]
        rs.alignment = Alignment(horizontal="center")

        ws3.cell(row=ri, column=8, value=entry["pii_tags"])
        ws3.cell(row=ri, column=9, value=entry["pii_count"])
        imgs = entry.get("embedded_images", 0)
        ws3.cell(row=ri, column=10, value=imgs if imgs > 0 else "-")
        ws3.cell(row=ri, column=11, value=entry["sha256"][:16] + "...")
        ws3.cell(row=ri, column=12, value=entry["last_modified"])

        for c in range(1, 13):
            _dcell(ws3.cell(row=ri, column=c), wrap=(c in (3, 8)))
        if ri % 2 == 0:
            for c in range(1, 13):
                if c not in (6, 7):
                    ws3.cell(row=ri, column=c).fill = GREY

    if file_audit:
        ws3.auto_filter.ref = f"A1:L{len(file_audit) + 1}"
    _autofit(ws3)

    # ================================================================
    #  TAB 4 — INDIAN REGULATORY GUIDE
    # ================================================================
    ws4 = wb.create_sheet("Indian Regulatory Guide")

    ws4.merge_cells("A1:G1")
    ws4["A1"] = ("Indian Data Protection Laws & Regulators — "
                 "PII Categories, Compliance & Remediation")
    ws4["A1"].font = SEC_F

    hdrs4 = [
        "PII Category", "Description", "Risk",
        "Primary Law / Regulator", "Secondary Rules",
        "Compliance Obligation", "Recommended Action"
    ]
    _hdr(ws4, 2, hdrs4)

    rules = [
        # HIGH RISK — Identity Documents
        ("AADHAAR",
         "Aadhaar Number (Verhoeff validated)",
         "HIGH",
         "DPDP Act 2023 Sec 2(t), Aadhaar Act 2016 Sec 29",
         "DPDP Rules 2025, UIDAI Circular on Masked Aadhaar",
         "Consent required. No public display. "
         "72-hr breach notification to DPBI.",
         "Encrypt at rest & transit. Use masked Aadhaar "
         "(last 4 digits only). Strict access control. "
         "Data localisation (store within India)."),

        ("PAN_CARD",
         "Permanent Account Number",
         "HIGH",
         "DPDP Act 2023, Income Tax Act Sec 139A",
         "DPDP Rules 2025, CBDT Notification",
         "Consent for collection. Erasure on request.",
         "Mask in reports/exports (show first 2 and last 1). "
         "Encrypt storage. Consent-based collection only."),

        ("PERSON_NAME",
         "Person Name (NER + Column Analysis)",
         "HIGH",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011 Sec 3, DPDP Rules 2025",
         "Lawful purpose & consent. Right to correction.",
         "Obtain explicit consent. Limit to need-to-know. "
         "Pseudonymise in analytics. Enable correction."),

        ("EMAIL",
         "Email Address",
         "HIGH",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011 Sec 3, DPDP Rules 2025",
         "Consent-based collection. Easy withdrawal.",
         "Consent-based collection. Anonymise where possible. "
         "Provide unsubscribe/withdrawal mechanism."),

        ("PHONE_IN",
         "Indian Mobile / Phone Number",
         "HIGH",
         "DPDP Act 2023, TRAI Privacy Regulations",
         "DPDP Rules 2025, TRAI DND Regulations",
         "Consent required. No sharing with telemarketers.",
         "Consent-based. Mask in logs/reports. Encrypt. "
         "Comply with TRAI DND preferences."),

        ("DATE_OF_BIRTH",
         "Date of Birth",
         "HIGH",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011 Sec 3, DPDP Rules 2025",
         "Minimise collection. Right to erasure.",
         "Collect only when legally required. Derive age "
         "where sufficient. Encrypt storage."),

        ("PASSPORT_IN",
         "Indian Passport Number",
         "HIGH",
         "DPDP Act 2023, Passports Act 1967",
         "DPDP Rules 2025",
         "Consent + purpose limitation.",
         "Encrypt. Restrict access to authorised personnel. "
         "Delete after purpose is fulfilled."),

        ("VOTER_ID",
         "Voter ID / EPIC Number",
         "HIGH",
         "DPDP Act 2023, RPA 1950",
         "DPDP Rules 2025",
         "Consent + purpose limitation.",
         "Encrypt. Restrict access. Audit usage."),

        ("DRIVING_LICENSE",
         "Indian Driving License Number",
         "HIGH",
         "DPDP Act 2023, Motor Vehicles Act 1988",
         "DPDP Rules 2025",
         "Consent + purpose limitation.",
         "Encrypt. Purpose limitation. Delete when no longer needed."),

        ("CREDIT_CARD",
         "Credit/Debit Card Number (Luhn validated)",
         "HIGH",
         "PCI-DSS v4.0",
         "RBI Card-on-File Tokenisation Directions 2022",
         "Never store full card number. Tokenise mandatorily.",
         "Tokenise immediately. Never store full PAN. "
         "Immediate remediation if found. PCI-DSS audit."),

        ("FINANCIAL",
         "Salary / CTC / Payroll Data",
         "HIGH",
         "DPDP Act 2023, Payment of Wages Act",
         "DPDP Rules 2025, Industrial Employment Act",
         "Consent for processing. Access restricted.",
         "Encrypt. Restrict to HR/Finance. Audit access logs. "
         "Retention per employment law requirements."),

        ("BANK_ACCOUNT",
         "Bank Account Number",
         "HIGH",
         "RBI Master Directions on IT Governance 2023",
         "DPDP Act 2023, DPDP Rules 2025",
         "Data localisation per RBI. Consent required.",
         "Encrypt. Mask in reports. Data localisation "
         "(store within India). Access audit."),

        ("HEALTH_DATA",
         "Medical / Health Information",
         "HIGH",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011 Sec 3 (sensitive personal data), "
         "DPDP Rules 2025",
         "Explicit consent. Higher protection standard.",
         "Encrypt. Restrict access. Explicit consent. "
         "Never share without authorisation. Data minimisation."),

        ("BIOMETRIC",
         "Biometric Data (fingerprint, iris, etc.)",
         "HIGH",
         "DPDP Act 2023, Aadhaar Act 2016 Sec 29",
         "DPDP Rules 2025, IT Rules 2011 Sec 3",
         "Explicit consent. No unauthorised storage.",
         "Encrypt. Store only if legally mandated. "
         "Restrict access to authorised systems only."),

        ("DEMAT_ACCOUNT",
         "Demat Account / DP ID",
         "HIGH",
         "SEBI CSCRF 2023",
         "DPDP Act 2023, SEBI Circular on KYC",
         "Protection mandated by SEBI for intermediaries.",
         "Encrypt. Access control. Audit trail required. "
         "Comply with SEBI CSCRF framework."),

        ("INSURANCE",
         "Insurance Policy Number",
         "HIGH",
         "IRDAI Info & Cybersecurity Guidelines 2023",
         "DPDP Act 2023, DPDP Rules 2025",
         "Consent required. Cybersecurity audit mandated.",
         "Encrypt. Restrict access. Comply with IRDAI "
         "cybersecurity framework. Regular audits."),

        ("PENSION",
         "PRAN / PF / UAN / Pension Data",
         "HIGH",
         "PFRDA Subscriber Data Protection Guidelines",
         "DPDP Act 2023, EPFO Rules, DPDP Rules 2025",
         "Protection of subscriber data mandated by PFRDA.",
         "Encrypt. Restrict to authorised personnel. "
         "Comply with PFRDA data protection guidelines."),

        # MEDIUM RISK
        ("EMBEDDED_IMAGES",
         "Embedded Images (may contain PII)",
         "MEDIUM",
         "DPDP Act 2023",
         "IT Rules 2011, DPDP Rules 2025",
         "Images may contain IDs, signatures, screenshots.",
         "Manual review required. Check for screenshots of IDs, "
         "signatures, scanned documents. OCR scanning recommended."),

        ("IFSC_CODE",
         "IFSC Code (Bank Branch)",
         "MEDIUM",
         "RBI Master Directions",
         "DPDP Act 2023",
         "Low risk alone; high when paired with account number.",
         "Protect when associated with account numbers. "
         "Contextual risk assessment needed."),

        ("GSTIN",
         "GST Identification Number",
         "MEDIUM",
         "GST Act 2017",
         "DPDP Act 2023",
         "Semi-public but protect when linked to personal data.",
         "Protect when linked to personal identity. "
         "Limited standalone risk."),

        ("UPI_ID",
         "UPI Virtual Payment Address",
         "MEDIUM",
         "RBI UPI Guidelines, NPCI Circulars",
         "DPDP Act 2023",
         "Consent-based usage. No unauthorised transactions.",
         "Consent-based. Mask in log files. Audit access."),

        ("EMPLOYEE_ID",
         "Employee Code / ID",
         "MEDIUM",
         "DPDP Act 2023",
         "Internal HR Policy, DPDP Rules 2025",
         "Restrict to authorised HR personnel.",
         "Access control. Audit usage."),

        ("DATE_EMPLOYMENT",
         "Date of Joining / Exit / Confirmation",
         "MEDIUM",
         "DPDP Act 2023",
         "Internal HR Policy",
         "Retain per retention policy. Delete on exit.",
         "Restrict access. Retain per company policy. "
         "Delete within retention period."),

        ("ADDRESS",
         "Physical / Postal Address",
         "MEDIUM",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011, DPDP Rules 2025",
         "Consent for collection. Right to correction.",
         "Minimise. Pseudonymise where possible. "
         "Encrypt when stored."),

        ("CLIENT_ID",
         "Client ID (SEBI regulated)",
         "MEDIUM",
         "SEBI CSCRF 2023",
         "DPDP Act 2023",
         "Protection as per SEBI cyber resilience framework.",
         "Access control. Audit. SEBI compliance."),

        ("FOLIO_NUMBER",
         "Mutual Fund Folio Number",
         "MEDIUM",
         "SEBI CSCRF 2023",
         "DPDP Act 2023",
         "Protected under SEBI investor data rules.",
         "Encrypt. Access control. SEBI compliance."),

        # LOW RISK
        ("ORGANISATION",
         "Organisation Name (NER)",
         "LOW",
         "DPDP Act 2023",
         "Companies Act 2013",
         "Non-sensitive. Protect if linked to individual data.",
         "Generally non-sensitive. Protect when linked to "
         "personal data of individual employees/clients."),

        ("LOCATION",
         "Location / City / State (NER)",
         "LOW",
         "DPDP Act 2023 Sec 2(t)",
         "IT Rules 2011",
         "Contextual risk with address/person data.",
         "Low risk alone. Higher when combined with "
         "address or person data."),

        ("IP_ADDRESS",
         "IPv4 Address",
         "LOW",
         "IT Act 2000 Sec 43",
         "CERT-In Directions 2022 (6-hr mandatory reporting)",
         "Log for security. Report incidents within 6 hours.",
         "Log for security purposes. Anonymise in analytics. "
         "CERT-In: Report cyber incidents within 6 hours."),

        ("PIN_CODE",
         "Indian PIN Code (6-digit)",
         "LOW",
         "DPDP Act 2023",
         "India Post",
         "Low risk alone. Contextual when with address.",
         "Low risk standalone. Higher when combined with "
         "full address data. Anonymise in analytics."),

        ("ORG_DATA",
         "Cost Center / Department / Designation",
         "LOW",
         "Internal Policy",
         "DPDP Act 2023",
         "Organisational metadata. Protect per policy.",
         "Internal data. Restrict per access policy. "
         "Low standalone risk."),
    ]

    for ri, (tag, desc, sens, primary, secondary, obligation,
             action) in enumerate(rules, 3):
        ws4.cell(row=ri, column=1, value=tag).font = B_FONT
        ws4.cell(row=ri, column=2, value=desc)
        rc = ws4.cell(row=ri, column=3, value=sens)
        if sens in RISK_STYLE:
            rc.fill, rc.font = RISK_STYLE[sens]
            rc.alignment = Alignment(horizontal="center")
        ws4.cell(row=ri, column=4, value=primary)
        ws4.cell(row=ri, column=5, value=secondary)
        ws4.cell(row=ri, column=6, value=obligation)
        ws4.cell(row=ri, column=7, value=action)
        for c in range(1, 8):
            _dcell(ws4.cell(row=ri, column=c), wrap=(c >= 2))
        ws4.cell(row=ri, column=1).font = B_FONT  # re-apply bold tag
        if ri % 2 == 0:
            for c in range(1, 8):
                if c != 3:
                    ws4.cell(row=ri, column=c).fill = GREY

    _autofit(ws4)

    # -- Save --
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    try:
        wb.save(output_path)
        print(f"\n  [OK] Report saved: {output_path}")
        return output_path
    except PermissionError:
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        dir_name = os.path.dirname(output_path)
        fallback = os.path.join(dir_name,
                                f"PII_Discovery_Report_{ts}.xlsx")
        wb.save(fallback)
        print("\n  [!] Report file is open in Excel.")
        print(f"  [OK] Saved to: {fallback}")
        return fallback


# ===================================================================
#  CLI ENTRY POINT
# ===================================================================

def main():
    parser = argparse.ArgumentParser(
        description=f"{TOOL_NAME} v{TOOL_VERSION} — India Regulatory Focus",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pii_scanner_india.py
  python pii_scanner_india.py --target "E:/Data/test"
  python pii_scanner_india.py --target "C:/Users/Docs" --output "D:/Reports"
        """)
    parser.add_argument("--target", "-t", type=str,
                        default=str(DEFAULT_TARGET),
                        help="Target folder to scan (default: ./test)")
    parser.add_argument("--output", "-o", type=str,
                        default=str(DEFAULT_REPORTS_DIR),
                        help="Output directory for reports (default: ./reports)")
    args = parser.parse_args()

    target = str(pathlib.Path(args.target).resolve())
    outdir = str(pathlib.Path(args.output).resolve())
    report_path = os.path.join(outdir, REPORT_FILENAME)

    nlp, model_name = get_nlp()
    ner_status = (f"ACTIVE ({model_name})"
                  if HAS_SPACY else "DISABLED")
    ocr_status = ("ACTIVE" if (HAS_PIL and HAS_TESSERACT)
                  else "NOT AVAILABLE")

    print()
    print("=" * 65)
    print(f"   {TOOL_NAME} v{TOOL_VERSION}")
    print("   India Regulatory Focus | Fully Offline | ML-Enhanced")
    print("=" * 65)
    print(f"  Target:       {target}")
    print(f"  Report:       {report_path}")
    print(f"  Scan Time:    "
          f"{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"  NER Engine:   {ner_status}")
    print(f"  OCR Engine:   {ocr_status}")
    print("  Regulations:  DPDP Act 2023 | DPDP Rules 2025 | IT Act 2000")
    print("                RBI | SEBI | IRDAI | PFRDA | TRAI | CERT-In")
    print("=" * 65)
    print()

    # Phase 1 — Scan
    print("[PHASE 1] Scanning files...")
    print()
    start_time = datetime.datetime.now()
    findings, file_audit = scan_folder(target)
    scan_duration = (datetime.datetime.now() - start_time).total_seconds()

    # Phase 2 — Report
    print()
    print("[PHASE 2] Generating Report...")
    generate_report(findings, file_audit, report_path, target)

    # Summary
    high_c = sum(1 for f in findings if f["sensitivity"] == "HIGH")
    med_c = sum(1 for f in findings if f["sensitivity"] == "MEDIUM")
    low_c = sum(1 for f in findings if f["sensitivity"] == "LOW")

    print()
    print("=" * 65)
    print("   SCAN COMPLETE")
    print("=" * 65)
    print(f"  Duration:           {scan_duration:.2f}s")
    print(f"  Files Scanned:      {len(file_audit)}")
    print(f"  Files with PII:     "
          f"{sum(1 for f in file_audit if f['status'] == 'PII_DETECTED')}")
    print(f"  Total Findings:     {len(findings)}")
    print(f"    HIGH Risk:        {high_c}")
    print(f"    MEDIUM Risk:      {med_c}")
    print(f"    LOW Risk:         {low_c}")
    print(f"  Report:             {report_path}")
    print("=" * 65)
    print()


if __name__ == "__main__":
    main()
